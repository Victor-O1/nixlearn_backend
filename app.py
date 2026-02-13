"""
Mr. Dexter - Master Teacher Backend
Compact FastAPI server with smart LLM-based chunking and RAG
"""
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional, List
import os, json, uuid, shutil
from pathlib import Path

from dotenv import load_dotenv
load_dotenv()
# Document processing
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation

# LLM & RAG
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_groq import ChatGroq
from langchain_core.prompts import ChatPromptTemplate
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.documents import Document

print("\n" + "="*80)
print("🎓 MR. DEXTER - MASTER TEACHER BACKEND")
print("="*80)

# Initialize FastAPI
app = FastAPI(title="Mr. Dexter API")
print("✅ [INIT] FastAPI app initialized")

# CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3002"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
print("✅ [INIT] CORS middleware configured")

# Configuration
GROQ_API_KEY = os.getenv("GROQ_API_KEY", "")
print(f"🔑 [INIT] Groq API key loaded: {'✅ YES' if GROQ_API_KEY else '❌ NO (SET GROQ_API_KEY!)'}")

UPLOAD_DIR = Path("./uploads")
SESSIONS_DIR = Path("./sessions")
UPLOAD_DIR.mkdir(exist_ok=True)
SESSIONS_DIR.mkdir(exist_ok=True)
print(f"📁 [INIT] Upload directory: {UPLOAD_DIR.absolute()}")
print(f"💾 [INIT] Sessions directory: {SESSIONS_DIR.absolute()}")

# Initialize LLM
print("🤖 [INIT] Initializing Groq LLM...")
llm = ChatGroq(
    groq_api_key=GROQ_API_KEY,
    model_name="llama-3.3-70b-versatile",
    temperature=0.3
)
print("✅ [INIT] Groq LLM ready (model: llama-3.3-70b-versatile, temp: 0.3)")

# Initialize embeddings
print("🧬 [INIT] Initializing embeddings model...")
embeddings = HuggingFaceEmbeddings(
    model_name="sentence-transformers/all-MiniLM-L6-v2"
)
print("✅ [INIT] Embeddings model ready (all-MiniLM-L6-v2)")

# Models
class SessionCreate(BaseModel):
    syllabus_text: Optional[str] = None

class ChunkResponse(BaseModel):
    chunk_id: str
    topic: str
    content: str
    summary: str
    order: int

class SessionResponse(BaseModel):
    session_id: str
    chunks: List[ChunkResponse]
    syllabus_topics: Optional[List[str]] = None

# ========== DOCUMENT EXTRACTION ==========

def extract_text_from_pdf(file_path: str) -> str:
    print(f"\n📄 [EXTRACT-PDF] Starting PDF extraction")
    print(f"📄 [EXTRACT-PDF] File: {file_path}")
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            total_pages = len(pdf.pages)
            print(f"📄 [EXTRACT-PDF] Total pages: {total_pages}")
            
            for i, page in enumerate(pdf.pages, 1):
                print(f"📄 [EXTRACT-PDF] Processing page {i}/{total_pages}...")
                page_text = page.extract_text() or ""
                text += page_text + "\n\n"
                print(f"📄 [EXTRACT-PDF] Page {i}: {len(page_text)} chars extracted")
                
        print(f"📄 [EXTRACT-PDF] ✅ Complete! Total text: {len(text)} chars")
        print(f"📄 [EXTRACT-PDF] Preview: {text[:200]}...")
        return text
    except Exception as e:
        print(f"📄 [EXTRACT-PDF] ❌ ERROR: {e}")
        raise

def extract_text_from_docx(file_path: str) -> str:
    print(f"\n📝 [EXTRACT-DOCX] Starting DOCX extraction")
    print(f"📝 [EXTRACT-DOCX] File: {file_path}")
    try:
        doc = DocxDocument(file_path)
        total_paras = len(doc.paragraphs)
        print(f"📝 [EXTRACT-DOCX] Total paragraphs: {total_paras}")
        
        text = "\n\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        print(f"📝 [EXTRACT-DOCX] ✅ Complete! Total text: {len(text)} chars")
        print(f"📝 [EXTRACT-DOCX] Preview: {text[:200]}...")
        return text
    except Exception as e:
        print(f"📝 [EXTRACT-DOCX] ❌ ERROR: {e}")
        raise

def extract_text_from_pptx(file_path: str) -> str:
    print(f"\n🎯 [EXTRACT-PPTX] Starting PPTX extraction")
    print(f"🎯 [EXTRACT-PPTX] File: {file_path}")
    try:
        prs = Presentation(file_path)
        total_slides = len(prs.slides)
        print(f"🎯 [EXTRACT-PPTX] Total slides: {total_slides}")
        
        text = ""
        for i, slide in enumerate(prs.slides, 1):
            print(f"🎯 [EXTRACT-PPTX] Processing slide {i}/{total_slides}...")
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            slide_content = "\n".join(slide_text)
            text += slide_content + "\n\n"
            print(f"🎯 [EXTRACT-PPTX] Slide {i}: {len(slide_content)} chars extracted")
            
        print(f"🎯 [EXTRACT-PPTX] ✅ Complete! Total text: {len(text)} chars")
        print(f"🎯 [EXTRACT-PPTX] Preview: {text[:200]}...")
        return text
    except Exception as e:
        print(f"🎯 [EXTRACT-PPTX] ❌ ERROR: {e}")
        raise

def extract_text_from_file(file_path: str, filename: str) -> str:
    print(f"\n🔍 [EXTRACT] Detecting file type for: {filename}")
    
    if filename.lower().endswith('.pdf'):
        print("🔍 [EXTRACT] Type detected: PDF")
        return extract_text_from_pdf(file_path)
    elif filename.lower().endswith('.docx'):
        print("🔍 [EXTRACT] Type detected: DOCX")
        return extract_text_from_docx(file_path)
    elif filename.lower().endswith('.pptx'):
        print("🔍 [EXTRACT] Type detected: PPTX")
        return extract_text_from_pptx(file_path)
    elif filename.lower().endswith('.txt'):
        print("🔍 [EXTRACT] Type detected: TXT")
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        print(f"🔍 [EXTRACT] TXT loaded: {len(text)} chars")
        return text
    else:
        print(f"🔍 [EXTRACT] ❌ Unsupported file type: {filename}")
        raise ValueError(f"Unsupported file type: {filename}")

# ========== SMART LLM-BASED CHUNKING ==========

def smart_topic_chunking(text: str, syllabus_topics: Optional[List[str]] = None) -> List[dict]:
    print("\n" + "="*80)
    print("🧠 [CHUNKING] STARTING SMART TOPIC-BASED CHUNKING")
    print("="*80)
    print(f"🧠 [CHUNKING] Input text length: {len(text)} chars")
    print(f"🧠 [CHUNKING] Syllabus provided: {'✅ YES' if syllabus_topics else '❌ NO'}")
    
    if syllabus_topics:
        print(f"🧠 [CHUNKING] Number of syllabus topics: {len(syllabus_topics)}")
        for i, topic in enumerate(syllabus_topics, 1):
            print(f"🧠 [CHUNKING]   {i}. {topic}")
        return chunk_by_syllabus(text, syllabus_topics)
    else:
        print("🧠 [CHUNKING] Mode: AUTO-DISCOVERY (LLM will find topics)")
        return chunk_by_auto_discovery(text)

def chunk_by_syllabus(text: str, syllabus_topics: List[str]) -> List[dict]:
    print(f"\n📚 [SYLLABUS-CHUNK] Syllabus-guided chunking started")
    print(f"📚 [SYLLABUS-CHUNK] Processing {len(syllabus_topics)} topics")
    
    topics_str = "\n".join([f"{i+1}. {t}" for i, t in enumerate(syllabus_topics)])
    print(f"📚 [SYLLABUS-CHUNK] Topics formatted:\n{topics_str}")
    
    matching_prompt = ChatPromptTemplate.from_messages([
        ("system", """You are a content analyzer. Given study material and syllabus topics, 
identify which sections of the material correspond to each topic.

Return a JSON array where each element has:
- topic: the syllabus topic name
- relevant_text: the extracted text section for this topic (can be multiple paragraphs)
- confidence: how confident you are this matches (0.0-1.0)

Only include topics that actually appear in the material. If a topic isn't covered, skip it."""),
        ("user", """Study Material:
{text}

Syllabus Topics:
{topics}

Return JSON only, no other text.""")
    ])
    
    print(f"📚 [SYLLABUS-CHUNK] Truncating text to 15000 chars for LLM...")
    text_truncated = text[:15000]
    print(f"📚 [SYLLABUS-CHUNK] Calling LLM for topic matching...")
    
    try:
        response = llm.invoke(matching_prompt.format_messages(text=text_truncated, topics=topics_str))
        print(f"📚 [SYLLABUS-CHUNK] ✅ LLM response received")
        print(f"📚 [SYLLABUS-CHUNK] Response length: {len(response.content)} chars")
        
        # Parse response
        content = response.content.strip()
        print(f"📚 [SYLLABUS-CHUNK] Parsing JSON response...")
        
        if content.startswith("```json"):
            print(f"📚 [SYLLABUS-CHUNK] Removing markdown code fence...")
            content = content[7:]
        if content.endswith("```"):
            content = content[:-3]
        
        chunks_data = json.loads(content.strip())
        print(f"📚 [SYLLABUS-CHUNK] ✅ JSON parsed successfully")
        print(f"📚 [SYLLABUS-CHUNK] Matched {len(chunks_data)} topics")
        
        chunks = []
        for i, chunk_data in enumerate(chunks_data):
            print(f"📚 [SYLLABUS-CHUNK] Processing matched topic {i+1}/{len(chunks_data)}...")
            chunk = {
                "chunk_id": str(uuid.uuid4()),
                "topic": chunk_data["topic"],
                "content": chunk_data["relevant_text"],
                "order": i
            }
            chunks.append(chunk)
            print(f"📚 [SYLLABUS-CHUNK]   ✅ Topic: {chunk['topic']}")
            print(f"📚 [SYLLABUS-CHUNK]   📊 Content length: {len(chunk['content'])} chars")
            print(f"📚 [SYLLABUS-CHUNK]   🆔 Chunk ID: {chunk['chunk_id']}")
        
        print(f"📚 [SYLLABUS-CHUNK] ✅ Syllabus chunking complete!")
        return chunks
        
    except Exception as e:
        print(f"📚 [SYLLABUS-CHUNK] ❌ ERROR: {e}")
        print(f"📚 [SYLLABUS-CHUNK] 🔄 FALLBACK: Using auto-discovery instead...")
        return chunk_by_auto_discovery(text)

def chunk_by_auto_discovery(text: str) -> List[dict]:
    print(f"\n🔍 [AUTO-CHUNK] Auto-discovery chunking started")
    
    discovery_prompt = ChatPromptTemplate.from_messages([
        ("system", """You are a content analyzer. Given study material, identify distinct topics/sections and split the content accordingly.

Return a JSON array where each element has:
- topic: a clear, descriptive topic name
- content: the text for this topic section
- order: numerical order (0, 1, 2, ...)

Look for natural topic boundaries like:
- Subject matter changes
- Conceptual shifts
- New chapters/sections
- Different themes

Aim for 5-15 meaningful chunks, not too granular."""),
        ("user", """Study Material:
{text}

Return JSON only, no other text.""")
    ])
    
    print(f"🔍 [AUTO-CHUNK] Truncating text to 15000 chars for LLM...")
    text_truncated = text[:15000]
    print(f"🔍 [AUTO-CHUNK] Calling LLM for topic discovery...")
    
    try:
        response = llm.invoke(discovery_prompt.format_messages(text=text_truncated))
        print(f"🔍 [AUTO-CHUNK] ✅ LLM response received")
        print(f"🔍 [AUTO-CHUNK] Response length: {len(response.content)} chars")
        
        # Parse response
        content = response.content.strip()
        print(f"🔍 [AUTO-CHUNK] Parsing JSON response...")
        
        if content.startswith("```json"):
            print(f"🔍 [AUTO-CHUNK] Removing markdown code fence...")
            content = content[7:]
        if content.endswith("```"):
            content = content[:-3]
        
        chunks_data = json.loads(content.strip())
        print(f"🔍 [AUTO-CHUNK] ✅ JSON parsed successfully")
        print(f"🔍 [AUTO-CHUNK] Discovered {len(chunks_data)} topics")
        
        chunks = []
        for chunk_data in chunks_data:
            chunk = {
                "chunk_id": str(uuid.uuid4()),
                "topic": chunk_data["topic"],
                "content": chunk_data["content"],
                "order": chunk_data.get("order", len(chunks))
            }
            chunks.append(chunk)
            print(f"🔍 [AUTO-CHUNK] Topic {chunk['order']}: {chunk['topic']}")
            print(f"🔍 [AUTO-CHUNK]   📊 Content length: {len(chunk['content'])} chars")
            print(f"🔍 [AUTO-CHUNK]   🆔 Chunk ID: {chunk['chunk_id']}")
        
        print(f"🔍 [AUTO-CHUNK] ✅ Auto-discovery chunking complete!")
        return chunks
        
    except Exception as e:
        print(f"🔍 [AUTO-CHUNK] ❌ ERROR: {e}")
        print(f"🔍 [AUTO-CHUNK] 🔄 FALLBACK: Using simple character-based splitting...")
        
        # Fallback: simple splitting
        splitter = RecursiveCharacterTextSplitter(chunk_size=2000, chunk_overlap=200)
        texts = splitter.split_text(text)
        print(f"🔍 [AUTO-CHUNK] Split into {len(texts)} chunks")
        
        chunks = []
        for i, chunk_text in enumerate(texts):
            chunks.append({
                "chunk_id": str(uuid.uuid4()),
                "topic": f"Section {i+1}",
                "content": chunk_text,
                "order": i
            })
            print(f"🔍 [AUTO-CHUNK] Fallback chunk {i+1}: {len(chunk_text)} chars")
        
        return chunks

# ========== SUMMARIZATION ==========

def generate_summary(chunk_content: str, topic: str) -> str:
    print(f"\n📝 [SUMMARY] Generating summary")
    print(f"📝 [SUMMARY] Topic: {topic}")
    print(f"📝 [SUMMARY] Content length: {len(chunk_content)} chars")
    
    summary_prompt = ChatPromptTemplate.from_messages([
        ("system", """You are Mr. Dexter, a master teacher. Create a clear, concise summary 
that helps students understand the key concepts.

Guidelines:
- 3-5 sentences max
- Focus on core concepts and key takeaways
- Use simple, clear language
- Help students grasp the essence quickly"""),
        ("user", """Topic: {topic}

Content:
{content}

Write a summary:""")
    ])
    
    print(f"📝 [SUMMARY] Truncating content to 3000 chars for LLM...")
    content_truncated = chunk_content[:3000]
    print(f"📝 [SUMMARY] Calling LLM for summary generation...")
    
    try:
        response = llm.invoke(summary_prompt.format_messages(topic=topic, content=content_truncated))
        summary = response.content.strip()
        print(f"📝 [SUMMARY] ✅ Summary generated: {len(summary)} chars")
        print(f"📝 [SUMMARY] Preview: {summary[:100]}...")
        return summary
    except Exception as e:
        print(f"📝 [SUMMARY] ❌ ERROR: {e}")
        print(f"📝 [SUMMARY] 🔄 FALLBACK: Using truncated content as summary")
        return chunk_content[:300] + "..."

# ========== RAG SETUP ==========

def setup_rag_for_session(session_id: str, chunks: List[dict]) -> Chroma:
    print(f"\n🗄️ [RAG] Setting up RAG vector store")
    print(f"🗄️ [RAG] Session ID: {session_id}")
    print(f"🗄️ [RAG] Number of chunks: {len(chunks)}")
    
    documents = []
    for i, chunk in enumerate(chunks, 1):
        print(f"🗄️ [RAG] Creating document {i}/{len(chunks)}...")
        doc = Document(
            page_content=chunk["content"],
            metadata={
                "chunk_id": chunk["chunk_id"],
                "topic": chunk["topic"],
                "order": chunk["order"],
                "summary": chunk.get("summary", "")
            }
        )
        documents.append(doc)
        print(f"🗄️ [RAG]   ✅ Document created for: {chunk['topic']}")
    
    persist_dir = str(SESSIONS_DIR / session_id / "vectorstore")
    print(f"🗄️ [RAG] Creating ChromaDB at: {persist_dir}")
    
    try:
        vectorstore = Chroma.from_documents(
            documents=documents,
            embedding=embeddings,
            persist_directory=persist_dir
        )
        print(f"🗄️ [RAG] ✅ Vector store created successfully!")
        print(f"🗄️ [RAG] Persist directory: {persist_dir}")
        return vectorstore
    except Exception as e:
        print(f"🗄️ [RAG] ❌ ERROR: {e}")
        raise

# ========== API ENDPOINTS ==========

@app.post("/api/sessions/create", response_model=SessionResponse)
async def create_session(
    files: List[UploadFile] = File(None),
    syllabus_text: Optional[str] = Form(None)
):
    print("\n" + "="*80)
    print("🎓 [SESSION] CREATE SESSION REQUEST RECEIVED")
    print("="*80)
    
    session_id = str(uuid.uuid4())
    print(f"🎓 [SESSION] Generated Session ID: {session_id}")
    
    session_dir = SESSIONS_DIR / session_id
    session_dir.mkdir(exist_ok=True)
    print(f"🎓 [SESSION] Created session directory: {session_dir}")
    
    # Parse syllabus
    syllabus_topics = None
    if syllabus_text and syllabus_text.strip():
        print(f"📚 [SYLLABUS] Syllabus text received: {len(syllabus_text)} chars")
        print(f"📚 [SYLLABUS] Raw syllabus:\n{syllabus_text}")
        print(f"📚 [SYLLABUS] Calling LLM to extract topics...")
        
        syllabus_prompt = ChatPromptTemplate.from_messages([
            ("system", "Extract a clean list of topics from this syllabus. Return JSON array of strings."),
            ("user", "{syllabus}\n\nReturn JSON only.")
        ])
        
        try:
            response = llm.invoke(syllabus_prompt.format_messages(syllabus=syllabus_text))
            content = response.content.strip()
            print(f"📚 [SYLLABUS] LLM response received: {len(content)} chars")
            
            if content.startswith("```json"):
                content = content[7:-3]
            syllabus_topics = json.loads(content)
            print(f"📚 [SYLLABUS] ✅ Extracted {len(syllabus_topics)} topics:")
            for i, topic in enumerate(syllabus_topics, 1):
                print(f"📚 [SYLLABUS]   {i}. {topic}")
        except Exception as e:
            print(f"📚 [SYLLABUS] ❌ LLM parsing failed: {e}")
            print(f"📚 [SYLLABUS] 🔄 FALLBACK: Using line-based splitting")
            syllabus_topics = [line.strip() for line in syllabus_text.split('\n') if line.strip()]
            print(f"📚 [SYLLABUS] Fallback extracted {len(syllabus_topics)} topics")
    else:
        print(f"📚 [SYLLABUS] No syllabus provided")
    
    # Process files
    combined_text = ""
    
    if files:
        print(f"📤 [UPLOAD] Processing {len(files)} uploaded file(s)")
        for i, file in enumerate(files, 1):
            if file.filename:
                print(f"📤 [UPLOAD] File {i}/{len(files)}: {file.filename}")
                file_path = UPLOAD_DIR / f"{session_id}_{file.filename}"
                
                print(f"📤 [UPLOAD] Saving to: {file_path}")
                with open(file_path, "wb") as buffer:
                    shutil.copyfileobj(file.file, buffer)
                print(f"📤 [UPLOAD] ✅ File saved")
                
                # Extract text
                text = extract_text_from_file(str(file_path), file.filename)
                combined_text += text + "\n\n"
                print(f"📤 [UPLOAD] Text extracted and added to combined text")
        
        print(f"📤 [UPLOAD] ✅ All files processed")
        print(f"📤 [UPLOAD] Total combined text: {len(combined_text)} chars")
    else:
        print(f"📤 [UPLOAD] No files uploaded")
    
    # Generate AI content if no files
    if not combined_text.strip():
        print("🤖 [AI-GEN] No content available, generating AI content...")
        if syllabus_topics:
            gen_prompt = f"Generate comprehensive study material for these topics:\n" + "\n".join(syllabus_topics)
            print(f"🤖 [AI-GEN] Using syllabus-based prompt")
        else:
            gen_prompt = "Generate sample educational content on a general topic."
            print(f"🤖 [AI-GEN] Using generic prompt")
        
        print(f"🤖 [AI-GEN] Calling LLM to generate content...")
        gen_response = llm.invoke(gen_prompt)
        combined_text = gen_response.content
        print(f"🤖 [AI-GEN] ✅ Generated {len(combined_text)} chars")
        print(f"🤖 [AI-GEN] Preview: {combined_text[:200]}...")
    
    # Smart chunking
    print(f"\n🔄 [PIPELINE] Starting chunking pipeline...")
    chunks = smart_topic_chunking(combined_text, syllabus_topics)
    print(f"🔄 [PIPELINE] ✅ Chunking complete: {len(chunks)} chunks created")
    
    # Generate summaries
    print(f"\n🔄 [PIPELINE] Starting summary generation for all chunks...")
    for i, chunk in enumerate(chunks, 1):
        print(f"🔄 [PIPELINE] Generating summary {i}/{len(chunks)}...")
        chunk["summary"] = generate_summary(chunk["content"], chunk["topic"])
    print(f"🔄 [PIPELINE] ✅ All summaries generated")
    
    # Setup RAG
    print(f"\n🔄 [PIPELINE] Setting up RAG...")
    vectorstore = setup_rag_for_session(session_id, chunks)
    print(f"🔄 [PIPELINE] ✅ RAG setup complete")
    
    # Save session metadata
    session_data = {
        "session_id": session_id,
        "syllabus_topics": syllabus_topics,
        "chunks": chunks
    }
    
    metadata_path = session_dir / "metadata.json"
    print(f"💾 [SAVE] Saving session metadata to: {metadata_path}")
    with open(metadata_path, "w") as f:
        json.dump(session_data, f, indent=2)
    print(f"💾 [SAVE] ✅ Metadata saved")
    
    print("\n" + "="*80)
    print("🎉 [SESSION] SESSION CREATED SUCCESSFULLY!")
    print(f"🎉 [SESSION] Session ID: {session_id}")
    print(f"🎉 [SESSION] Total chunks: {len(chunks)}")
    print(f"🎉 [SESSION] Syllabus topics: {len(syllabus_topics) if syllabus_topics else 0}")
    print("="*80 + "\n")
    
    return SessionResponse(
        session_id=session_id,
        chunks=[ChunkResponse(**chunk) for chunk in chunks],
        syllabus_topics=syllabus_topics
    )

@app.get("/api/sessions/{session_id}", response_model=SessionResponse)
async def get_session(session_id: str):
    print(f"\n🔍 [GET] Retrieving session: {session_id}")
    
    session_dir = SESSIONS_DIR / session_id
    metadata_file = session_dir / "metadata.json"
    
    print(f"🔍 [GET] Looking for metadata at: {metadata_file}")
    
    if not metadata_file.exists():
        print(f"🔍 [GET] ❌ Session not found!")
        raise HTTPException(status_code=404, detail="Session not found")
    
    print(f"🔍 [GET] ✅ Metadata file found, loading...")
    with open(metadata_file, "r") as f:
        session_data = json.load(f)
    
    print(f"🔍 [GET] ✅ Session loaded successfully")
    print(f"🔍 [GET] Chunks: {len(session_data['chunks'])}")
    print(f"🔍 [GET] Syllabus topics: {len(session_data.get('syllabus_topics', [])) if session_data.get('syllabus_topics') else 0}")
    
    return SessionResponse(
        session_id=session_data["session_id"],
        chunks=[ChunkResponse(**chunk) for chunk in session_data["chunks"]],
        syllabus_topics=session_data.get("syllabus_topics")
    )

@app.get("/api/health")
async def health_check():
    print("❤️ [HEALTH] Health check requested")
    print("❤️ [HEALTH] Status: healthy")
    return {"status": "healthy", "service": "Mr. Dexter API"}

print("\n✅ All endpoints registered")
print("📍 POST /api/sessions/create - Create new session")
print("📍 GET /api/sessions/{session_id} - Get session")
print("📍 GET /api/health - Health check")

if __name__ == "__main__":
    import uvicorn
    print("\n" + "="*80)
    print("🚀 STARTING MR. DEXTER API SERVER")
    print("="*80)
    print("📍 Server URL: http://localhost:8000")
    print("📚 API Docs: http://localhost:8000/docs")
    print("🔧 Interactive API: http://localhost:8000/redoc")
    print("="*80 + "\n")
    uvicorn.run(app, host="0.0.0.0", port=8000)